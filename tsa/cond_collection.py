#!/usr/bin/python
# -*- coding: utf-8 -*-

# Collection of Conditions for analysis

import logging
import pptx
import os
import openpyxl as xl
from .condition import Condition
from .error import TsaErrCollection
from .utils import strfdelta
from .utils import list_local_statids
from .utils import list_local_sensors
from collections import OrderedDict
from datetime import datetime
from io import BytesIO
from pptx.util import Pt
from pptx.util import Cm
from pptx.dml.color import RGBColor
from openpyxl.cell.cell import WriteOnlyCell

log = logging.getLogger(__name__)

class CondCollection:
    """
    A collection of conditions to analyze.
    All conditions share the same analysis time range.
    Times are assumed to be given as dates only,
    and their HH:MM:SS are overridden to default values,
    see ``set_default_times(self)``.

    :param time_from: start time (inclusive) of the analysis period
    :type time_from: datetime
    :param time_until: end time (exclusive) of the analysis period
    :type time_until: datetime
    :param title: name of the collection
    :type title: str
    """
    def __init__(self, time_from, time_until, title):
        # Times must be datetime objects and in correct order
        assert isinstance(time_from, datetime)
        assert isinstance(time_until, datetime)
        assert time_from <= time_until
        self.time_from = time_from.replace(hour=0, minute=0, second=0)
        self.time_until = time_until.replace(hour=23, minute=59, second=59)

        self.title = title
        # Timestamp is based on instance creation time,
        # not when the analysis has been run
        self.created_at = datetime.now()

        # Containers for conditions and unique stations in them.
        # Combinations of site and master_alias are used as unique
        # identifiers among the conditions.
        self.conditions = OrderedDict()

        # Database-specific stuff
        self.has_main_db_view = False
        self.station_ids_in_db_view = set()

        self.errors = TsaErrCollection(f'COLLECTION <{self.title}>')

    def add_condition(self, site, master_alias, raw_condition, excel_row=None):
        """
        Add new Condition instance, except if one exists already
        with same site-master_alias identifier.
        """
        try:
            candidate = Condition(
                site=site,
                master_alias=master_alias,
                raw_condition=raw_condition,
                time_range=(self.time_from, self.time_until),
                excel_row=excel_row
                )
        except:
            self.errors.add(
                msg=f'Could not build Condition, skipping (Excel row {excel_row})',
                log_add='exception'
            )
            return
        if candidate.id_string in self.conditions.keys():
            self.errors.add(
                msg=f'Condition identifier "{candidate.id_string}" is already reserved, skipping (Excel row {excel_row})',
                log_add='warning'
            )
            return
        self.conditions[candidate.id_string] = candidate

    def setup_obs_view(self, pg_conn):
        """
        Create temporary view ``obs_main``
        that works as the main source for Block queries.

        :param pg_conn: valid psycopg2 connection object
        """
        from_str = self.time_from.strftime('%Y-%m-%d %H:%M:%S')
        until_str = self.time_until.strftime('%Y-%m-%d %H:%M:%S')
        sql = ("CREATE OR REPLACE TEMP VIEW obs_main AS "
               "SELECT tfrom, statid, seid, seval "
               "FROM statobs "
               "INNER JOIN seobs "
               "ON statobs.id = seobs.obsid "
               f"WHERE tfrom BETWEEN '{from_str}'::timestamptz AND '{until_str}'::timestamptz;")
        with pg_conn.cursor() as cur:
            try:
                log.debug(sql)
                cur.execute(sql)
                pg_conn.commit()
                self.has_main_db_view = True
            except:
                pg_conn.rollback()
                self.errors.add(msg='Cannot create obs_main db view',
                                log_add='exception')

    def validate_statids_with_db(self, pg_conn):
        """
        Fetch unique station ids from main obs db view,
        and check for each primary Block if its station id
        exists in the view.
        If not, record an error to the block.

        .. note: Unlike statid validation with non-db set,
            this db validation is done at CondCollection level,
            not at AnalysisCollection level.
            In AnalysisCollection, the validation stationid set is always
            the same. When querying unique stationids from database,
            the result can differ between CondCollections if they apply
            different time range limits for the ``main_obs`` view.
        """
        # FIXME: This query takes too much time at the moment
        sql = "SELECT DISTINCT statid FROM obs_main ORDER BY statid;"
        with pg_conn.cursor() as cur:
            try:
                log.info('Fetching unique station ids from db view `obs_main` ...')
                cur.execute(sql)
                statids_from_db = cur.fetchall()
                statids_from_db = set(el[0] for el in statids_from_db)
            except:
                self.errors.add(msg=('Cannot fetch station ids for Block validation '
                                     'from db view obs_main'),
                                log_add='exception')
                return
        for c in self.conditions.keys():
            for b in self.conditions[c].blocks.keys():
                log.info(('Db stationid validation for '
                           f'{str(self.conditions[c].blocks[b])} of '
                           f'{str(self.conditions[c])} of {str(self)} ...'))
                isprimary = self.conditions[c].blocks[b].secondary is False
                hasid = self.conditions[c].blocks[b].station_id is not None
                validstatid = self.conditions[c].blocks[b].station_id in statids_from_db
                if not isprimary:
                    continue
                if not hasid:
                    self.conditions[c].blocks[b].errors.add(
                        msg='stationid is None (tried to compare it to ids from db view)',
                        log_add='error'
                    )
                    continue
                if not validstatid:
                    self.conditions[c].blocks[b].errors.add(
                        msg='stationid was not found in ids from db view',
                        log_add='error'
                    )

    def create_condition_temptables(self, pg_conn):
        """
        For each Condition, create the corresponding temporary table in db.
        Primary conditions are handled first, only then secondary ones;
        if there are secondary conditions depending further on each other,
        it is up to the user to give them in correct order!
        """
        # First round for primary ones only
        # so temp tables referenced by secondary conditions
        # can be found in the database session
        for cnd in self.conditions.keys():
            if self.conditions[cnd].secondary or not self.conditions[cnd].is_valid():
                continue
            self.conditions[cnd].create_db_temptable(pg_conn=pg_conn)

        # Second round for secondary ones,
        # viewnames list is now updated every time
        for cnd in self.conditions.keys():
            if not self.conditions[cnd].is_valid():
                continue
            if self.conditions[cnd].secondary:
                self.conditions[cnd].create_db_temptable(pg_conn=pg_conn)

    def fetch_all_results(self, pg_conn):
        """
        Fetch results
        for all Conditions that have a corresponding view in the database.
        """
        cnd_len = len(self.conditions)
        for i, cnd in enumerate(self.conditions.keys()):
            log.info(f'Fetching {i+1}/{cnd_len}: {str(self.conditions[cnd])} ...')
            try:
                self.conditions[cnd].fetch_results_from_db(pg_conn=pg_conn)
            except:
                self.conditions[cnd].errors.add(
                    msg='Exception while fetching results, skipping',
                    log_add='exception'
                )

    def as_text(self, value):
        return str(value) if value is not None else ""

    def adjust_column_width(self, worksheet):
        for column_cells in worksheet.columns:
            length = max(len(self.as_text(cell.value)) for cell in column_cells)
            worksheet.column_dimensions[xl.utils.get_column_letter(column_cells[0].column)].width = length

    def to_worksheet_per_site(self, cnd, wb):
        """
        Add a worksheet to an ``openpyxl.Workbook`` instance
        containing row based results of the condition collection.
        """
        df = cnd.main_df
        ws2 = wb.create_sheet()
        ws2.title = cnd.site

        cell = WriteOnlyCell(ws2)
        cell.style = 'Pandas'

        # print(f"### DEBUG, site: {cnd.site}")
        columns = list(df.columns.values)
        # print(f"### DEBUG, columns: {columns}")
        for col_num in range(len(columns)):
            ws2.cell(row=1, column=col_num+1).value = str(columns[col_num])

        row_count = df.shape[0]
        # print(f"### DEBUG, row_count: {row_count}")
        sheet_row = 2
        for row_num in range(row_count):
            sheet_row = row_num+2
            for col_num in range(len(columns)):
                # print(f"### cond_collection DEBUG, col_num: {col_num}")
                # print(f"### cond_collection DEBUG, str(df[columns[[{col_num}]][{i}]): {str(df[columns[col_num]][i])}")
                ws2.cell(row=sheet_row, column=col_num+1).value = str(df[columns[col_num]][row_num])

        self.adjust_column_width(ws2)

    def to_worksheet(self, wb):
        """
        Add a worksheet to an ``openpyxl.Workbook`` instance
        containing summary results of the condition collection.
        """
        assert isinstance(wb, xl.Workbook)
        ws = wb.create_sheet()
        ws.title = self.title or 'conditions'

        # Headers in fixed cells & styling
        headers = {'A1': 'start',
                   'B1': 'end',
                   'D1': 'analyzed',
                   'A3': 'site',
                   'B3': 'master_alias',
                   'C3': 'condition',
                   'D3': 'data_from',
                   'E3': 'data_until',
                   'F3': 'valid',
                   'G3': 'notvalid',
                   'H3': 'nodata',
                   'I3': 'rows'
                   }

        for k, v in headers.items():
            ws[k] = v
            ws[k].font = xl.styles.Font(bold=True)

        # Global values
        ws['A2'] = self.time_from
        ws['B2'] = self.time_until
        ws['D2'] = self.created_at

        # Condition rows
        r = 4
        for cnd in self.conditions.values():
            ws[f'A{r}'] = cnd.site
            ws[f'B{r}'] = cnd.master_alias
            ws[f'C{r}'] = cnd.condition
            ws[f'D{r}'] = cnd.data_from
            ws[f'E{r}'] = cnd.data_until
            ws[f'F{r}'] = cnd.percentage_valid
            ws[f'G{r}'] = cnd.percentage_notvalid
            ws[f'H{r}'] = cnd.percentage_nodata
            ws[f'I{r}'] = cnd.main_df.shape[0]

            # Percent format
            ws[f'F{r}'].number_format = '0.00 %'
            ws[f'G{r}'].number_format = '0.00 %'
            ws[f'H{r}'].number_format = '0.00 %'

            # Add new sheets for row based data of analysis
            self.to_worksheet_per_site(cnd, wb)

            r += 1

        self.adjust_column_width(ws)

    def to_pptx(self, pptx_template, png_dir=None):
        """
        Return a ``pptx`` presentation object,
        making a slide of each condition.

        ``pptx`` must be a filepath or file-like object
        representing a PowerPoint file that includes the master
        layout for the TSA report and nothing else. The default
        placeholder indices must conform with the constants here!
        """
        phi = dict(
        HEADER_IDX = 17,     # Slide header placeholder
        TITLE_IDX = 0,       # Condition title placeholder
        BODY_IDX = 13,       # Condition string placeholder
        TIMERANGE_IDX = 15,  # Placeholder for condition start/end time text
        VALIDTABLE_IDX = 18, # Validity time/percentage table placeholder
        ERRORS_IDX = 19,     # Placeholder for errors and warnings
        MAINPLOT_IDX = 11,   # Main timeline plot placeholder
        FOOTER_IDX = 16,     # Slide footer placeholder
        )
        MAINPLOT_H_PX = 3840 # Main timeline plot height in pixels

        pres = pptx.Presentation(pptx_template)
        layout = pres.slide_layouts[0]

        # Ensure placeholder indices exist as they should
        indices_in_pres = [ph.placeholder_format.idx for ph in layout.placeholders]
        for k, v in phi.items():
            if v not in indices_in_pres:
                raise Exception(f'{k} {v} not in default layout placeholders')

        # Add slides and fill in contents for each condition.
        for c in self.conditions.values():
            s = pres.slides.add_slide(layout)

            # Slide header
            txt = 'TSA report: '
            if self.title is not None:
                txt += self.title
            txt += ' ' + self.created_at.strftime('%d.%m.%Y')
            s.placeholders[phi['HEADER_IDX']].text = txt

            # Slide footer
            txt = 'TSATool v0.1, copyright WSP Finland'
            s.placeholders[phi['FOOTER_IDX']].text = txt

            # Condition title
            s.placeholders[phi['TITLE_IDX']].text = c.id_string

            # Condition string / body
            s.placeholders[phi['BODY_IDX']].text = c.condition

            # Condition data time range
            if not (c.data_from is None or c.data_until is None):
                txt = 'Datan tarkasteluväli {}-{}'.format(
                    c.data_from.strftime('%d.%m.%Y %H:%M'),
                    c.data_until.strftime('%d.%m.%Y %H:%M')
                )
            else:
                txt = 'Ei dataa saatavilla'
            s.placeholders[phi['TIMERANGE_IDX']].text = txt

            # Master condition validity table
            tb_shape = s.placeholders[phi['VALIDTABLE_IDX']].insert_table(rows=3, cols=4)
            tb = tb_shape.table

            tb.cell(0, 0).text = ''

            tb.cell(0, 1).text = 'Voimassa'
            tb.cell(0, 2).text = 'Ei voimassa'
            tb.cell(0, 3).text = 'Tieto puuttuu'

            tb.cell(1, 0).text = 'Yhteensä'
            txt = strfdelta(c.tottime_valid, '{days} pv {hours} h {minutes} min')
            tb.cell(1, 1).text = txt
            txt = strfdelta(c.tottime_notvalid, '{days} pv {hours} h {minutes} min')
            tb.cell(1, 2).text = txt
            txt = strfdelta(c.tottime_nodata, '{days} pv {hours} h {minutes} min')
            tb.cell(1, 3).text = txt

            tb.cell(2, 0).text = 'Osuus tarkasteluajasta'
            txt = '{} %'.format(round(c.percentage_valid*100, 2))
            tb.cell(2, 1).text = txt
            txt = '{} %'.format(round(c.percentage_notvalid*100, 2))
            tb.cell(2, 2).text = txt
            txt = '{} %'.format(round(c.percentage_nodata*100, 2))
            tb.cell(2, 3).text = txt

            for cl in tb.iter_cells():
                cl.fill.background()
                for ph in cl.text_frame.paragraphs:
                    ph.font.name = 'Montserrat'
                    ph.font.size = Pt(8)
                    ph.font.color.rgb = RGBColor.from_string('000000')

            for row in tb.rows:
                row.height = Cm(0.64)

            # Condition errors and warnings
            txt = c.errors.short_str()
            s.placeholders[phi['ERRORS_IDX']].text = txt

            # Condition main timeline plot; ignored if no data to viz
            if c.main_df is None:
                continue
            # Find out the proportion of plot height of the width
            wh_factor = s.placeholders[phi['MAINPLOT_IDX']].height \
                        / s.placeholders[phi['MAINPLOT_IDX']].width
            w, h = MAINPLOT_H_PX, wh_factor*MAINPLOT_H_PX

            # NOTE: Saving png as in-memory object does not work
            #       for some reason.
            #       Saving it to file instead and keeping the file
            #       if png_dir is provided.
            # with BytesIO() as fobj:
            #     c.save_timelineplot(fobj, w, h)
            #     s.placeholders[phi['MAINPLOT_IDX']].insert_picture(fobj)
            if png_dir is None:
                fobj = 'temp.png'
                rm_png = True
            else:
                if not os.path.exists(png_dir):
                    self.errors.add(
                        msg=f'Directory "{png_dir}" for images does not exist, not saving png files',
                        log_add='warning'
                    )
                    fobj = 'temp.png'
                    rm_png = True
                else:
                    fobj = os.path.join(png_dir, f'{self.title}_{c.id_string}.png')
                    rm_png = False
            saved = c.save_timelineplot(fobj, w, h)
            if saved:
                s.placeholders[phi['MAINPLOT_IDX']].insert_picture(fobj)
                if rm_png:
                    os.remove(fobj)

        return pres

    def save_pptx(self, pptx_template, out_path, png_dir=None):
        """
        Call ``.to_pptx`` and save result to file.
        """
        pptx_obj = self.to_pptx(pptx_template=pptx_template, png_dir=png_dir)
        pptx_obj.save(out_path)

    def run_analysis(self,
                     pg_conn,
                     wb=None,
                     wb_path=None,
                     pptx_path=None,
                     pptx_template=None,
                     png_dir=None):
        """
        Call necessary methods to run the condition analysis
        and save results to the specified
        ``openpyxl.Workbook`` instance ``wb`` as new worksheet
        and the ``pptx_path`` as ``.pptx`` file.
        If ``wb_path`` is provided, save the workbook in the end
        (will overwrite existing files).
        If an output is ``None``, it is not created.
        """
        log.info(f'Starting analysis of {str(self)}')
        self.setup_obs_view(pg_conn=pg_conn)
        log.info('obs_main db view created')
        # FIXME: Station id validation agains unique values in db view
        #        is not done, because the SELECT DISTINCT query is very
        #        slow for some reason.
        #        This step is not mandatory, though.
        #        If a station id is missing, the result should be just
        #        an empty table and / or a database error for that condition.
        # self.validate_statids_with_db(pg_conn=pg_conn)
        # log.debug('Station ids validated')
        self.create_condition_temptables(pg_conn=pg_conn)
        log.info('Temp tables created for conditions')

        log.info('Starting to fetch results from database ...')
        starttime = datetime.now()
        self.fetch_all_results(pg_conn=pg_conn)
        log.info(f'Results fetched in {str(datetime.now() - starttime)}')

        if wb is not None:
            log.info('Creating Excel sheet ...')
            self.to_worksheet(wb)
            if wb_path is not None:
                wb.save(wb_path)
                log.info(f'Excel sheet saved to {wb_path}')
        else:
            log.warning(f'No Excel sheet saved from {str(self)}')

        if pptx_path is not None and pptx_template is not None:
            log.info(f'Saving Powerpoint report as {pptx_path} ...')
            self.save_pptx(pptx_template=pptx_template,
                           out_path=pptx_path,
                           png_dir=png_dir)
            log.info(f'{pptx_path} saved')
        else:
            log.warning(f'No Powerpoint report saved from {str(self)}')

    def __getitem__(self, key):
        """
        Returns the Condition instance on the corresponding index.
        """
        return self.conditions[key]

    def __str__(self):
        s = (f'<CondCollection {self.title} '
             f'with {len(self.conditions)} conditions>')
        return s

    @classmethod
    def from_xlsx_sheet(cls, ws):
        """
        Create a condition collection for analysis
        based on an ``openpyxl`` ``worksheet`` object ``ws``.

        .. note:: Start and end dates must be in cells A2 and B2, respectively,
                  and conditions must be listed starting from row 4,
                  such that ``site`` is in column A,
                  ``master_alias`` in column B
                  and ``raw_condition`` in column C.
                  There must not be empty rows in between.
                  Any columns outside A:C are ignored,
                  so additional data can be placed outside them.
        """
        # Start and end dates must be d.m.Y dates, start in cell A2
        # and end in cell B2.
        dateformat = '%d.%m.%Y'
        time_from, time_until = ws['A2'].value, ws['B2'].value
        if time_from is None:
            raise Exception('Start date in cell A2 is empty')
        if not isinstance(time_from, datetime):
            try:
                time_from = datetime.strptime(time_from, dateformat)
            except:
                raise Exception('Cannot parse start date in cell A2')
        if time_until is None:
            raise Exception('End date in cell B2 is empty')
        if not isinstance(time_until, datetime):
            try:
                time_until = datetime.strptime(time_until, dateformat)
            except:
                raise Exception('Cannot parse end date in cell B2')
        if time_from > time_until:
            raise Exception('Start date in cell A2 must not be greater than end date in cell B2')

        cc = cls(time_from=time_from, time_until=time_until, title=ws.title)
        for row in ws.iter_rows(min_row=4, max_col=3):
            cells = [c for c in row]
            cells_ok = True
            for c in cells:
                if c.value is None:
                    cc.errors.add(f'Cell {c.coordinate} is empty: condition row ignored')
                    cells_ok = False
            # Row is ignored if any of the three cells is empty
            if not cells_ok:
                continue
            cc.add_condition(site=cells[0].value, master_alias=cells[1].value,
                             raw_condition=cells[2].value, excel_row=cells[0].row)

        return cc
